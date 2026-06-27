#!/usr/bin/env python3
"""Build Photosynthesis.pptx — 6-slide biology presentation with green nature palette."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────
DARK_GREEN   = RGBColor(0x1B, 0x43, 0x32)   # title / conclusion bg
LEAF_GREEN   = RGBColor(0x2D, 0x6A, 0x4F)   # accent
MID_GREEN    = RGBColor(0x40, 0x91, 0x6C)   # lighter accent
HIGHLIGHT    = RGBColor(0x52, 0xB7, 0x88)   # bright highlight
LIGHT_BG     = RGBColor(0xF0, 0xF7, 0xF4)   # content slide bg
OFF_WHITE    = RGBColor(0xFA, 0xFD, 0xFB)   # near-white
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT    = RGBColor(0x0D, 0x28, 0x18)   # near-black dark
BODY_TEXT    = RGBColor(0x1B, 0x43, 0x32)   # body on light
SUBTLE       = RGBColor(0x6B, 0x90, 0x80)   # muted green-gray
WARM_ACCENT  = RGBColor(0xD4, 0xA3, 0x37)   # gold accent for highlights
GOLD_LIGHT   = RGBColor(0xF0, 0xD0, 0x60)   # lighter gold

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape

def rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape

def circle_shape(slide, left, top, size, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape

def arrow_right(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def arrow_down(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def text_box(slide, left, top, width, height, text="", font_name="Calibri",
             font_size=Pt(18), bold=False, color=BODY_TEXT, alignment=PP_ALIGN.LEFT,
             line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.line_spacing = line_spacing
    return txBox

def add_paragraph(tf, text, font_name="Calibri", font_size=Pt(18), bold=False,
                  color=BODY_TEXT, alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                  line_spacing=1.15, space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    p.line_spacing = line_spacing
    p.space_after = space_after
    return p

def multi_text(slide, left, top, width, height, paragraphs_data, default_size=Pt(18),
               default_color=BODY_TEXT, default_font="Calibri"):
    """paragraphs_data: list of (text, bold?, font_size|None, color|None)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(paragraphs_data):
        if isinstance(item, str):
            text, bold, fs, clr = item, False, default_size, default_color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            fs = item[2] if len(item) > 2 else default_size
            clr = item[3] if len(item) > 3 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = default_font
        p.font.size = fs
        p.font.bold = bold
        p.font.color.rgb = clr
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
    return txBox


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide (Levy, presenter)
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_GREEN)

# Decorative top accent bar
rect_shape(slide1, Inches(0), Inches(0), W, Inches(0.08), HIGHLIGHT)

# Large decorative leaf strip (subtle geometric shapes)
for i in range(5):
    x = Inches(0.5 + i * 2.6)
    circle_shape(slide1, x, Inches(2.0), Inches(0.7), MID_GREEN)

# Main title
text_box(slide1, Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.5),
         "How Plants Eat Sunshine",
         font_name="Calibri Light", font_size=Pt(54), bold=False, color=WHITE,
         alignment=PP_ALIGN.CENTER)

text_box(slide1, Inches(1.5), Inches(3.65), Inches(10.3), Inches(1.0),
         "The Magic of Photosynthesis",
         font_name="Calibri Light", font_size=Pt(48), bold=False, color=GOLD_LIGHT,
         alignment=PP_ALIGN.CENTER)

# Separator line
rect_shape(slide1, Inches(4.5), Inches(4.55), Inches(4.3), Inches(0.03), HIGHLIGHT)

# Subtitle
text_box(slide1, Inches(1.5), Inches(4.85), Inches(10.3), Inches(0.7),
         "Or: Why You Owe Every Breath to a Leaf",
         font_name="Calibri Light", font_size=Pt(26), bold=False, color=SUBTLE,
         alignment=PP_ALIGN.CENTER)

# Team line
text_box(slide1, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.5),
         "Levy  •  Member 2  •  Member 3  •  Member 4  •  Member 5",
         font_name="Calibri", font_size=Pt(18), bold=False, color=SUBTLE,
         alignment=PP_ALIGN.CENTER)

# Bottom accent bar
rect_shape(slide1, Inches(0), Inches(7.42), W, Inches(0.08), HIGHLIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — What IS Photosynthesis? (Member 2)
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, LIGHT_BG)

# Top accent
rect_shape(slide2, Inches(0), Inches(0), W, Inches(0.06), DARK_GREEN)

# Section number + title
txBox_s2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(5), Inches(0.8))
tf = txBox_s2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "01  ·  WHAT IS PHOTOSYNTHESIS?"
p.font.name = "Calibri"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = LEAF_GREEN
p.line_spacing = 1.0

text_box(slide2, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.7),
         "Plants Turn Sunlight into Food",
         font_name="Calibri Light", font_size=Pt(36), bold=False, color=DARK_GREEN)

# Equation box
eq_box = rounded_rect(slide2, Inches(0.8), Inches(1.85), Inches(7.5), Inches(1.25),
                      fill_color=RGBColor(0xE0, 0xEE, 0xE8),
                      line_color=LEAF_GREEN, line_width=Pt(1.5))

txBox_eq = slide2.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(7.0), Inches(1.1))
tf_eq = txBox_eq.text_frame
tf_eq.word_wrap = True
p_eq = tf_eq.paragraphs[0]
p_eq.text = "The Simplest Definition"
p_eq.font.name = "Calibri"
p_eq.font.size = Pt(14)
p_eq.font.bold = True
p_eq.font.color.rgb = LEAF_GREEN
p_eq2 = tf_eq.add_paragraph()
p_eq2.text = "Sunlight  +  Water  +  CO₂    ⟶    Glucose  +  Oxygen"
p_eq2.font.name = "Calibri"
p_eq2.font.size = Pt(20)
p_eq2.font.bold = False
p_eq2.font.color.rgb = DARK_TEXT
p_eq2.space_before = Pt(4)

# Equation formula below
text_box(slide2, Inches(0.8), Inches(3.2), Inches(7.5), Inches(0.5),
         "6 CO₂  +  6 H₂O  +  Light Energy    ⟶    C₆H₁₂O₆  +  6 O₂",
         font_name="Calibri", font_size=Pt(16), bold=False, color=DARK_TEXT)

# Fun fact callout
fun_box = rounded_rect(slide2, Inches(8.8), Inches(1.85), Inches(3.7), Inches(1.95),
                       fill_color=DARK_GREEN, line_color=None)
txBox_fun = slide2.shapes.add_textbox(Inches(9.0), Inches(1.95), Inches(3.3), Inches(1.75))
tf_fun = txBox_fun.text_frame
tf_fun.word_wrap = True
p_fun = tf_fun.paragraphs[0]
p_fun.text = "💡  Fun Fact"
p_fun.font.name = "Calibri"
p_fun.font.size = Pt(14)
p_fun.font.bold = True
p_fun.font.color.rgb = GOLD_LIGHT
p_fun2 = tf_fun.add_paragraph()
p_fun2.text = "Every oxygen molecule you breathe was once split from a water molecule by a plant."
p_fun2.font.name = "Calibri Light"
p_fun2.font.size = Pt(14)
p_fun2.font.color.rgb = WHITE
p_fun2.space_before = Pt(8)
p_fun2.line_spacing = 1.3

# Where it happens section
text_box(slide2, Inches(0.8), Inches(4.0), Inches(7.0), Inches(0.5),
         "Where It Happens",
         font_name="Calibri", font_size=Pt(16), bold=True, color=LEAF_GREEN)

text_box(slide2, Inches(0.8), Inches(4.35), Inches(7.0), Inches(0.7),
         "Inside chloroplasts — tiny green solar panels in leaf cells",
         font_name="Calibri Light", font_size=Pt(16), bold=False, color=BODY_TEXT)

# 3 inputs → 2 outputs visual
# Inputs row
y_io = Inches(5.3)
for idx, (label, clr) in enumerate([
    ("☀  SUNLIGHT", GOLD_LIGHT),
    ("💧  WATER", RGBColor(0x74, 0xBF, 0xE1)),
    ("🌬  CO₂", RGBColor(0xB0, 0xBE, 0xC5)),
]):
    bx = Inches(0.8 + idx * 2.5)
    btn = rounded_rect(slide2, bx, y_io, Inches(2.1), Inches(0.55),
                       fill_color=DARK_GREEN)
    text_box(slide2, Inches(0.95 + idx * 2.5), Inches(5.35), Inches(1.9), Inches(0.45),
             label, font_name="Calibri", font_size=Pt(13), bold=True, color=clr,
             alignment=PP_ALIGN.CENTER)

# Arrow
arrow_right(slide2, Inches(7.5), Inches(5.35), Inches(0.6), Inches(0.45), HIGHLIGHT)

# Outputs
for idx, (label, clr) in enumerate([
    ("🍬  GLUCOSE (Food)", GOLD_LIGHT),
    ("🫧  OXYGEN", RGBColor(0x74, 0xBF, 0xE1)),
]):
    bx = Inches(8.5 + idx * 2.2)
    btn = rounded_rect(slide2, bx, y_io, Inches(1.9), Inches(0.55),
                       fill_color=DARK_GREEN)
    text_box(slide2, Inches(8.6 + idx * 2.2), Inches(5.35), Inches(1.7), Inches(0.45),
             label, font_name="Calibri", font_size=Pt(13), bold=True, color=clr,
             alignment=PP_ALIGN.CENTER)

# Bottom accent
rect_shape(slide2, Inches(0), Inches(7.44), W, Inches(0.06), DARK_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — The Light Reactions (Member 3)
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, LIGHT_BG)
rect_shape(slide3, Inches(0), Inches(0), W, Inches(0.06), DARK_GREEN)

# Section label
text_box(slide3, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "02  ·  THE LIGHT REACTIONS", font_name="Calibri", font_size=Pt(13),
         bold=True, color=LEAF_GREEN)

text_box(slide3, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.7),
         "Catching Energy — Phase 1",
         font_name="Calibri Light", font_size=Pt(36), bold=False, color=DARK_GREEN)

text_box(slide3, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.5),
         'Happens in the thylakoid membranes — the "pancake stacks" inside chloroplasts',
         font_name="Calibri Light", font_size=Pt(16), bold=False, color=BODY_TEXT)

# ── Process flow visual ──
# Sun / Light source
circle_shape(slide3, Inches(1.2), Inches(2.6), Inches(1.2),
             fill_color=GOLD_LIGHT, line_color=RGBColor(0xE0, 0xA8, 0x20), line_width=Pt(2))
text_box(slide3, Inches(1.2), Inches(3.9), Inches(1.2), Inches(0.4),
         "Light Energy", font_name="Calibri", font_size=Pt(11), bold=True,
         color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Arrow sun → chlorophyll
arrow_right(slide3, Inches(2.6), Inches(3.0), Inches(0.55), Inches(0.35), HIGHLIGHT)

# Chlorophyll / thylakoid box
chloro = rounded_rect(slide3, Inches(3.4), Inches(2.3), Inches(2.8), Inches(1.8),
                      fill_color=DARK_GREEN, line_color=None)
text_box(slide3, Inches(3.55), Inches(2.4), Inches(2.5), Inches(0.45),
         "CHLOROPHYLL", font_name="Calibri", font_size=Pt(15), bold=True,
         color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)
text_box(slide3, Inches(3.55), Inches(2.8), Inches(2.5), Inches(0.9),
         "Captures light\nSplits H₂O molecules\nReleases O₂",
         font_name="Calibri Light", font_size=Pt(12), bold=False,
         color=WHITE, alignment=PP_ALIGN.CENTER)

# Arrow chlorophyll → products
arrow_right(slide3, Inches(6.5), Inches(3.0), Inches(0.55), Inches(0.35), HIGHLIGHT)

# Products boxes
y_prod = Inches(2.2)
for idx, (label, desc, clr) in enumerate([
    ("ATP", "Energy\nCurrency", GOLD_LIGHT),
    ("NADPH", "Electron\nCarrier", RGBColor(0x95, 0xD5, 0xB2)),
    ("O₂", "Oxygen\nReleased", RGBColor(0x74, 0xBF, 0xE1)),
]):
    bx = rounded_rect(slide3, Inches(7.35), Inches(2.15 + idx * 1.05),
                      Inches(1.5), Inches(0.9),
                      fill_color=DARK_GREEN if idx < 2 else MID_GREEN,
                      line_color=None)
    text_box(slide3, Inches(7.45), Inches(2.2 + idx * 1.05), Inches(1.3), Inches(0.3),
             label, font_name="Calibri", font_size=Pt(14), bold=True,
             color=clr, alignment=PP_ALIGN.CENTER)
    text_box(slide3, Inches(7.45), Inches(2.45 + idx * 1.05), Inches(1.3), Inches(0.5),
             desc, font_name="Calibri Light", font_size=Pt(10), bold=False,
             color=WHITE if idx < 2 else DARK_GREEN, alignment=PP_ALIGN.CENTER)

# Key points below
pts3 = [
    ("Chlorophyll captures sunlight and splits water", False),
    ("Produces ATP — the cell's energy currency", False),
    ("Produces NADPH — carries electrons to Phase 2", False),
    ("Byproduct: oxygen is released into the air", False),
]
y_pts = Inches(4.7)
for idx, (txt, bld) in enumerate(pts3):
    dot = circle_shape(slide3, Inches(0.9), Inches(4.82 + idx * 0.38), Inches(0.14), fill_color=LEAF_GREEN)
    text_box(slide3, Inches(1.2), Inches(4.78 + idx * 0.38), Inches(8), Inches(0.35),
             txt, font_name="Calibri Light", font_size=Pt(15), bold=False, color=BODY_TEXT)

# Analogy callout
analogy3 = rounded_rect(slide3, Inches(9.3), Inches(4.7), Inches(3.2), Inches(1.3),
                        fill_color=RGBColor(0xE0, 0xEE, 0xE8), line_color=LEAF_GREEN, line_width=Pt(1))
text_box(slide3, Inches(9.5), Inches(4.8), Inches(2.8), Inches(0.35),
         "🔋  Analogy", font_name="Calibri", font_size=Pt(13), bold=True, color=LEAF_GREEN)
text_box(slide3, Inches(9.5), Inches(5.15), Inches(2.8), Inches(0.7),
         'This is the "charging the battery" phase — capturing and storing energy.',
         font_name="Calibri Light", font_size=Pt(12), bold=False, color=BODY_TEXT)

# Bottom note
text_box(slide3, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
         "⚡  Light reactions require sunlight — they stop in the dark.",
         font_name="Calibri", font_size=Pt(12), bold=False, color=SUBTLE)

rect_shape(slide3, Inches(0), Inches(7.44), W, Inches(0.06), DARK_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — The Calvin Cycle (Member 4)
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, LIGHT_BG)
rect_shape(slide4, Inches(0), Inches(0), W, Inches(0.06), DARK_GREEN)

# Section label
text_box(slide4, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "03  ·  THE CALVIN CYCLE", font_name="Calibri", font_size=Pt(13),
         bold=True, color=LEAF_GREEN)

text_box(slide4, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.7),
         "Building Sugar — Phase 2",
         font_name="Calibri Light", font_size=Pt(36), bold=False, color=DARK_GREEN)

text_box(slide4, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.5),
         "Happens in the stroma — the liquid-filled space inside the chloroplast",
         font_name="Calibri Light", font_size=Pt(16), bold=False, color=BODY_TEXT)

# ── Calvin Cycle Diagram ──
# Center of cycle
cx, cy = Inches(4.2), Inches(4.2)
cycle_r = Inches(1.2)
cycle = circle_shape(slide4, cx - cycle_r, cy - cycle_r, cycle_r * 2,
                     fill_color=None, line_color=LEAF_GREEN, line_width=Pt(3))

text_box(slide4, Inches(3.45), Inches(3.95), Inches(1.5), Inches(0.5),
         "CALVIN\nCYCLE", font_name="Calibri", font_size=Pt(14), bold=True,
         color=DARK_GREEN, alignment=PP_ALIGN.CENTER)

# Inputs around the cycle
inputs_cycle = [
    # (label, angle-deg, distance from center)
    ("CO₂ from air", 0, Inches(2.2)),
    ("ATP from\nLight Rxns", 120, Inches(2.4)),
    ("NADPH from\nLight Rxns", 240, Inches(2.4)),
]
import math
for label, deg, dist in inputs_cycle:
    rad = math.radians(deg - 90)  # -90 so 0 degrees is top
    bx = cx + int(dist * math.cos(rad)) - Inches(0.65)
    by = cy + int(dist * math.sin(rad)) - Inches(0.4)
    inp_box = rounded_rect(slide4, bx, by, Inches(1.3), Inches(0.7),
                           fill_color=RGBColor(0xE0, 0xEE, 0xE8),
                           line_color=LEAF_GREEN, line_width=Pt(1))
    text_box(slide4, bx, by, Inches(1.3), Inches(0.6),
             label, font_name="Calibri", font_size=Pt(11), bold=True,
             color=DARK_GREEN, alignment=PP_ALIGN.CENTER)

# Output arrow on right
arrow_right(slide4, Inches(5.8), Inches(3.95), Inches(0.5), Inches(0.35), HIGHLIGHT)

# Glucose output
glu_box = rounded_rect(slide4, Inches(6.6), Inches(3.5), Inches(1.8), Inches(1.1),
                       fill_color=DARK_GREEN, line_color=None)
text_box(slide4, Inches(6.7), Inches(3.6), Inches(1.6), Inches(0.4),
         "GLUCOSE", font_name="Calibri", font_size=Pt(15), bold=True,
         color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)
text_box(slide4, Inches(6.7), Inches(3.95), Inches(1.6), Inches(0.5),
         "C₆H₁₂O₆\nPlant Food!",
         font_name="Calibri Light", font_size=Pt(12), bold=False,
         color=WHITE, alignment=PP_ALIGN.CENTER)

# Key points on left side
pts4 = [
    ("Uses CO₂ from the air + ATP & NADPH from light reactions", False),
    ("Builds glucose through a cycle of enzyme reactions", False),
    ("Named after Melvin Calvin — discovered in the 1950s", False),
    ("This phase doesn't need light directly!", True),
]
y_p4 = Inches(5.5)
for idx, (txt, bld) in enumerate(pts4):
    if bld:
        dot = arrow_right(slide4, Inches(0.7), Inches(5.6 + idx * 0.38), Inches(0.25), Inches(0.2), WARM_ACCENT)
    else:
        dot = circle_shape(slide4, Inches(0.75), Inches(5.65 + idx * 0.38), Inches(0.12), fill_color=LEAF_GREEN)
    text_box(slide4, Inches(1.1), Inches(5.52 + idx * 0.38), Inches(7.5), Inches(0.35),
             txt, font_name="Calibri Light", font_size=Pt(15), bold=bld,
             color=DARK_GREEN if bld else BODY_TEXT)

# Analogy callout
analogy4 = rounded_rect(slide4, Inches(9.0), Inches(5.5), Inches(3.5), Inches(1.3),
                        fill_color=RGBColor(0xE0, 0xEE, 0xE8), line_color=LEAF_GREEN, line_width=Pt(1))
text_box(slide4, Inches(9.2), Inches(5.6), Inches(3.1), Inches(0.35),
         "🏗️  Analogy", font_name="Calibri", font_size=Pt(13), bold=True, color=LEAF_GREEN)
text_box(slide4, Inches(9.2), Inches(5.9), Inches(3.1), Inches(0.8),
         'The "spending the battery to build stuff" phase — using stored energy to make sugar.',
         font_name="Calibri Light", font_size=Pt(12), bold=False, color=BODY_TEXT)

rect_shape(slide4, Inches(0), Inches(7.44), W, Inches(0.06), DARK_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Why This Matters to YOU (Member 5)
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, LIGHT_BG)
rect_shape(slide5, Inches(0), Inches(0), W, Inches(0.06), DARK_GREEN)

text_box(slide5, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "04  ·  WHY THIS MATTERS", font_name="Calibri", font_size=Pt(13),
         bold=True, color=LEAF_GREEN)

text_box(slide5, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.7),
         "Photosynthesis = Life Support for Planet Earth",
         font_name="Calibri Light", font_size=Pt(36), bold=False, color=DARK_GREEN)

# Four impact cards in a 2x2 grid
card_data = [
    ("OXYGEN", "Every second breath you take comes from ocean phytoplankton — not just trees.",
     "🫧"),
    ("FOOD CHAIN", "No photosynthesis = no plants = no animals = no humans. It's the base of everything.",
     "🌿"),
    ("CLIMATE", "Plants absorb CO₂, the main greenhouse gas. They're Earth's natural thermostat.",
     "🌍"),
    ("DEFORESTATION", "Cutting down forests is literally turning off Earth's oxygen factories.",
     "🪓"),
]

card_w = Inches(5.6)
card_h = Inches(1.35)
start_x = Inches(0.8)
start_y = Inches(1.95)
gap_x = Inches(0.3)
gap_y = Inches(0.25)

for idx, (title, desc, emoji) in enumerate(card_data):
    col = idx % 2
    row = idx // 2
    cx = start_x + col * (card_w + gap_x)
    cy = start_y + row * (card_h + gap_y)

    card = rounded_rect(slide5, cx, cy, card_w, card_h,
                        fill_color=WHITE, line_color=RGBColor(0xD0, 0xE4, 0xD8), line_width=Pt(1))

    # Left accent bar
    rect_shape(slide5, cx, cy, Inches(0.08), card_h,
               fill_color=LEAF_GREEN if idx != 3 else WARM_ACCENT)

    # Title
    text_box(slide5, cx + Inches(0.3), cy + Inches(0.12), Inches(4.8), Inches(0.35),
             title, font_name="Calibri", font_size=Pt(16), bold=True,
             color=LEAF_GREEN if idx != 3 else WARM_ACCENT)

    # Description
    text_box(slide5, cx + Inches(0.3), cy + Inches(0.55), Inches(4.8), Inches(0.7),
             desc, font_name="Calibri Light", font_size=Pt(14), bold=False,
             color=BODY_TEXT)

# Mind-blowing stat at bottom
stat_box = rounded_rect(slide5, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.0),
                        fill_color=DARK_GREEN, line_color=None)
text_box(slide5, Inches(1.2), Inches(5.25), Inches(11.0), Inches(0.7),
         "🌳  A single large tree produces enough oxygen for 2–4 humans per year.",
         font_name="Calibri Light", font_size=Pt(20), bold=False, color=GOLD_LIGHT,
         alignment=PP_ALIGN.CENTER)

# Bottom note
text_box(slide5, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
         "The ocean's phytoplankton produce 50–80% of Earth's oxygen. Every other breath = ocean.",
         font_name="Calibri", font_size=Pt(12), bold=False, color=SUBTLE,
         alignment=PP_ALIGN.CENTER)

rect_shape(slide5, Inches(0), Inches(7.44), W, Inches(0.06), DARK_GREEN)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Key Takeaways + Q&A (Levy)
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, DARK_GREEN)
rect_shape(slide6, Inches(0), Inches(0), W, Inches(0.08), HIGHLIGHT)

text_box(slide6, Inches(1.2), Inches(0.9), Inches(10.9), Inches(0.7),
         "Key Takeaways",
         font_name="Calibri Light", font_size=Pt(44), bold=False, color=WHITE,
         alignment=PP_ALIGN.CENTER)

# Separator
rect_shape(slide6, Inches(5.0), Inches(1.55), Inches(3.3), Inches(0.03), HIGHLIGHT)

# Takeaways
takeaways = [
    ("1", "Photosynthesis turns sunlight into chemical energy —\nthe most important chemical reaction on Earth."),
    ("2", "Two phases: Light Reactions (capture energy)\n→  Calvin Cycle (build sugar)."),
    ("3", "Every food chain on Earth starts with photosynthesis."),
    ("4", "Protecting forests = protecting the planet's life support system."),
]

for idx, (num, txt) in enumerate(takeaways):
    y = Inches(2.15 + idx * 1.05)
    # Number circle
    circle_shape(slide6, Inches(1.4), y, Inches(0.55), fill_color=HIGHLIGHT)
    text_box(slide6, Inches(1.4), Inches(y / 914400 + 0.1), Inches(0.55), Inches(0.55),
             num, font_name="Calibri", font_size=Pt(20), bold=True, color=DARK_GREEN,
             alignment=PP_ALIGN.CENTER)

    # Text
    text_box(slide6, Inches(2.3), y, Inches(9.5), Inches(1.0),
             txt, font_name="Calibri Light", font_size=Pt(18), bold=False, color=WHITE,
             alignment=PP_ALIGN.LEFT)

# Closing quote box
quote_box = rounded_rect(slide6, Inches(1.4), Inches(6.35), Inches(10.5), Inches(0.8),
                         fill_color=RGBColor(0x24, 0x52, 0x3F), line_color=HIGHLIGHT, line_width=Pt(1))
text_box(slide6, Inches(1.8), Inches(6.45), Inches(9.7), Inches(0.6),
         '"Next time you see a leaf — that\'s a solar-powered food factory. Respect it."',
         font_name="Calibri Light", font_size=Pt(20), bold=False, color=GOLD_LIGHT,
         alignment=PP_ALIGN.CENTER)

# Q&A
text_box(slide6, Inches(1.2), Inches(7.15), Inches(10.9), Inches(0.3),
         "Questions & Answers",
         font_name="Calibri", font_size=Pt(14), bold=True, color=SUBTLE,
         alignment=PP_ALIGN.CENTER)

rect_shape(slide6, Inches(0), Inches(7.42), W, Inches(0.08), HIGHLIGHT)

# ── Save ─────────────────────────────────────────────────────
out_path = r"C:\Users\LENOVO\projects\biology-pptx\Photosynthesis.pptx"
prs.save(out_path)
print(f"✅ Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: {os.path.getsize(out_path):,} bytes")
